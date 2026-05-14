#!/usr/bin/env python3
"""NotebookLM PPT Watermark Remover

Remove the NotebookLM watermark (logo + text) from the bottom-right corner
of slides in exported PPTX files.
CHỈ CẦN CHẠY TRỰC TIẾP (CLICK ĐÚP) ĐỂ TỰ ĐỘNG XỬ LÝ TẤT CẢ FILE PPTX TRONG THƯ MỤC
    python3 remove_notebooklm_watermark.py
"""

import argparse
import io
import sys
from pathlib import Path

from PIL import Image, ImageFilter
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- KIỂM TRA THƯ VIỆN XỬ LÝ ẢNH NÂNG CAO ---
try:
    import cv2
    import numpy as np
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

# Default watermark region parameters (offset from bottom-right corner)
DEFAULT_WM_WIDTH = 175
DEFAULT_WM_HEIGHT = 30
DEFAULT_MARGIN_RIGHT = 3
DEFAULT_MARGIN_BOTTOM = 3
SAMPLE_PADDING = 5

# Minimum coverage ratio to consider a picture as a full-slide background
MIN_COVERAGE = 0.85

# Feather width in pixels for edge blending (dành cho fallback cách cũ)
FEATHER_PX = 6

# XML namespaces
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
NS_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def has_watermark(img, wm_left, wm_top, wm_right, wm_bottom):
    """Detect whether a watermark likely exists in the given region."""
    w, h = img.size

    # Sample a background reference strip above the watermark
    ref_top = max(0, wm_top - (wm_bottom - wm_top) - SAMPLE_PADDING)
    ref_bottom = wm_top
    ref_strip = img.crop((wm_left, ref_top, wm_right, ref_bottom))
    wm_strip = img.crop((wm_left, wm_top, wm_right, wm_bottom))

    ref_pixels = list(ref_strip.getdata())
    wm_pixels = list(wm_strip.getdata())

    if not ref_pixels or not wm_pixels:
        return False

    def avg_color(pixels):
        channels = len(pixels[0]) if isinstance(pixels[0], tuple) else 1
        if channels == 1:
            return (sum(pixels) / len(pixels),)
        sums = [0] * channels
        for px in pixels:
            for c in range(channels):
                sums[c] += px[c]
        return tuple(s / len(pixels) for s in sums)

    ref_avg = avg_color(ref_pixels)
    wm_avg = avg_color(wm_pixels)

    # Compare average color difference across RGB channels
    diff = sum(abs(r - w) for r, w in zip(ref_avg[:3], wm_avg[:3])) / 3
    return diff > 3.0


def remove_watermark(img, wm_width, wm_height, margin_right, margin_bottom, perfect_mask=None):
    """Remove the NotebookLM watermark from the bottom-right corner."""
    w, h = img.size

    wm_left = w - wm_width - margin_right
    wm_top = h - wm_height - margin_bottom
    wm_right = w - margin_right
    wm_bottom = h - margin_bottom

    # Check if watermark exists
    if not has_watermark(img, wm_left, wm_top, wm_right, wm_bottom):
        return None

    if HAS_CV2:
        img_np = np.array(img)
        is_rgba = (len(img_np.shape) == 3 and img_np.shape[2] == 4)

        if is_rgba:
            bgr_img = cv2.cvtColor(img_np, cv2.COLOR_RGBA2BGR)
        else:
            bgr_img = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)

        roi = bgr_img[wm_top:wm_bottom, wm_left:wm_right]

        if perfect_mask is not None and perfect_mask.shape == roi.shape[:2]:
            # Đã có Smart Mask hoàn hảo từ phân tích đa slide
            dilated_mask = perfect_mask
        else:
            # Dự phòng: Slide đơn (Hạ ngưỡng Canny để bắt được logo mờ)
            gray_roi = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
            
            # Ngưỡng trung bình để bắt đủ nét mảnh của logo
            edges = cv2.Canny(gray_roi, 50, 150)
            
            # Kernel vừa đủ ôm sát logo, không quá to
            kernel = np.ones((3, 3), np.uint8)
            dilated_mask = cv2.dilate(edges, kernel, iterations=1)

        full_mask = np.zeros(bgr_img.shape[:2], dtype=np.uint8)
        full_mask[wm_top:wm_bottom, wm_left:wm_right] = dilated_mask

        # InpaintRadius = 3 để hút đủ mảng màu bù vào khoảng trống mà không bị nhòe
        # Dùng NS (Navier-Stokes) để giữ thẳng các đường biên nét chữ xung quanh
        inpainted = cv2.inpaint(bgr_img, full_mask, inpaintRadius=3, flags=cv2.INPAINT_NS)

        if is_rgba:
            b, g, r = cv2.split(inpainted)
            a = img_np[:, :, 3]
            final_cv = cv2.merge((r, g, b, a))
        else:
            final_cv = cv2.cvtColor(inpainted, cv2.COLOR_BGR2RGB)

        return Image.fromarray(final_cv)

    else:
        # ---------------------------------------------------------
        # CÁCH CŨ: CLONE TOOL (DÁN ĐÈ CƠ BẢN) DÀNH CHO MÁY CHƯA CÀI OPENCV
        # ---------------------------------------------------------
        result = img.copy()
        
        sample_top = max(0, wm_top - wm_height - SAMPLE_PADDING)
        sample_bottom = wm_top

        sample_strip = result.crop((wm_left, sample_top, wm_right, sample_bottom))

        target_w = wm_right - wm_left
        target_h = wm_bottom - wm_top
        sample_strip = sample_strip.resize((target_w, target_h), Image.LANCZOS)

        # Create a feathered alpha mask for smooth blending
        mask = Image.new("L", (target_w, target_h), 255)
        for y in range(min(FEATHER_PX, target_h)):
            alpha = int(255 * y / FEATHER_PX)
            for x in range(target_w):
                mask.putpixel((x, y), alpha)

        # Also feather left edge
        for x in range(min(FEATHER_PX, target_w)):
            alpha_x = int(255 * x / FEATHER_PX)
            for y in range(target_h):
                current = mask.getpixel((x, y))
                mask.putpixel((x, y), min(current, alpha_x))

        # Ensure sample_strip matches result mode for compositing
        if result.mode == "RGBA" and sample_strip.mode != "RGBA":
            sample_strip = sample_strip.convert("RGBA")
        elif result.mode == "RGB" and sample_strip.mode != "RGB":
            sample_strip = sample_strip.convert("RGB")

        result.paste(sample_strip, (wm_left, wm_top), mask)

        return result


def is_fullpage_image(shape, slide_width, slide_height):
    """Check if a picture shape covers most of the slide area."""
    w_ratio = shape.width / slide_width
    h_ratio = shape.height / slide_height
    return w_ratio >= MIN_COVERAGE and h_ratio >= MIN_COVERAGE


def process_pptx(input_path, output_path, wm_width, wm_height,
                 margin_right, margin_bottom):
    """Process a single PPTX file. Returns (processed, skipped) counts."""
    prs = Presentation(str(input_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    
    # ----------------------------------------------------------------
    # BƯỚC 1: QUÉT TRƯỚC ĐỂ THU THẬP TẤT CẢ CÁC ẢNH VÀ HÌNH NỀN
    # ----------------------------------------------------------------
    unique_images = []
    processed_part_ids = set()

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            if not is_fullpage_image(shape, slide_w, slide_h):
                continue

            try:
                pic = shape._element
                blip = pic.find(f".//{NS_A}blip")
                if blip is None:
                    continue
                r_id = blip.get(f"{NS_R}embed")
                if r_id is None:
                    continue
                image_part = slide.part.rels[r_id].target_part
            except Exception:
                continue

            part_id = id(image_part)
            if part_id in processed_part_ids:
                continue
            processed_part_ids.add(part_id)

            try:
                img_blob = shape.image.blob
                content_type = shape.image.content_type
                img = Image.open(io.BytesIO(img_blob))
                unique_images.append({
                    'part': image_part,
                    'img': img,
                    'content_type': content_type
                })
            except Exception:
                continue

    if not unique_images:
        prs.save(str(output_path))
        return 0, 0

    # ----------------------------------------------------------------
    # BƯỚC 2: PHÂN TÍCH ĐA SLIDE (MULTI-SLIDE ANALYSIS) ĐỂ TẠO MẶT NẠ CHUẨN
    # ----------------------------------------------------------------
    perfect_mask = None
    if HAS_CV2 and len(unique_images) >= 3:
        rois = []
        for item in unique_images:
            img = item['img']
            w, h = img.size
            wm_left = w - wm_width - margin_right
            wm_top = h - wm_height - margin_bottom
            wm_right = w - margin_right
            wm_bottom = h - margin_bottom
            
            if wm_left < 0 or wm_top < 0:
                continue
                
            img_np = np.array(img.convert('RGB'))
            roi = img_np[wm_top:wm_bottom, wm_left:wm_right]
            
            # Lọc ra các slide có hình nền thực sự khác nhau
            is_unique = True
            for existing_roi in rois:
                diff = np.mean(np.abs(existing_roi.astype(int) - roi.astype(int)))
                if diff < 5.0: # Background giống hệt slide trước
                    is_unique = False
                    break
            if is_unique:
                rois.append(roi)
        
        # Nếu có từ 3 slide background khác nhau trở lên, logo là thứ duy nhất KHÔNG ĐỔI
        if len(rois) >= 3:
            stack = np.stack(rois)
            # Tính độ lệch chuẩn của từng pixel trên tất cả các slide
            std = np.std(stack, axis=0)
            std_gray = np.mean(std, axis=2)
            
            # Lấy những pixel cố định (logo nằm yên)
            _, mask = cv2.threshold(std_gray.astype(np.uint8), 8, 255, cv2.THRESH_BINARY_INV)
            
            # ĐÃ BỎ BƯỚC DỌN NHIỄU MORPH_OPEN ĐỂ KHÔNG XÓA NHẦM NÉT CHỮ MẢNH CỦA LOGO
            
            # Kernel giãn nở thu nhỏ tối đa: Chỉ ôm sát viền mờ của logo
            kernel_dilate = np.ones((3, 3), np.uint8)
            perfect_mask = cv2.dilate(mask, kernel_dilate, iterations=1)

    # ----------------------------------------------------------------
    # BƯỚC 3: TIẾN HÀNH XÓA WATERMARK VÀ LƯU LẠI
    # ----------------------------------------------------------------
    processed = 0
    skipped = 0
    
    for item in unique_images:
        img = item['img']
        image_part = item['part']
        content_type = item['content_type']
        
        cleaned = remove_watermark(img, wm_width, wm_height,
                                   margin_right, margin_bottom, perfect_mask)
        if cleaned is None:
            skipped += 1
            continue

        # Save with correct format, preserving alpha
        buf = io.BytesIO()
        if "png" in content_type:
            fmt = "PNG"
        elif img.mode == "RGBA":
            fmt = "PNG"
        else:
            fmt = "JPEG"
        cleaned.save(buf, format=fmt)
        buf.seek(0)

        image_part._blob = buf.read()
        processed += 1

    prs.save(str(output_path))
    return processed, skipped


def main():
    if not HAS_CV2:
        print("💡 MẸO: Code đang dùng thuật toán dán đè cơ bản (nhìn hơi giống clone tool).", file=sys.stderr)
        print("💡 ĐỂ XÓA MƯỢT VÀ TỰ NHIÊN HƠN: Hãy mở Terminal/CMD và chạy lệnh: pip install opencv-python numpy\n", file=sys.stderr)

    parser = argparse.ArgumentParser(
        description="Remove NotebookLM watermark from PPTX slides"
    )
    # Changed nargs to "*" so it allows 0 arguments (for 1-click execution)
    parser.add_argument(
        "input", nargs="*", type=Path, help="Input PPTX file(s) (để trống sẽ tự động tìm trong thư mục hiện tại)"
    )
    parser.add_argument(
        "-o", "--output", type=Path, default="clean.pptx",
        help="Output file path (only for single file input)"
    )
    parser.add_argument(
        "--wm-width", type=int, default=DEFAULT_WM_WIDTH,
        help=f"Watermark width in pixels (default: {DEFAULT_WM_WIDTH})"
    )
    parser.add_argument(
        "--wm-height", type=int, default=DEFAULT_WM_HEIGHT,
        help=f"Watermark height in pixels (default: {DEFAULT_WM_HEIGHT})"
    )
    parser.add_argument(
        "--margin-right", type=int, default=DEFAULT_MARGIN_RIGHT,
        help=f"Right margin in pixels (default: {DEFAULT_MARGIN_RIGHT})"
    )
    parser.add_argument(
        "--margin-bottom", type=int, default=DEFAULT_MARGIN_BOTTOM,
        help=f"Bottom margin in pixels (default: {DEFAULT_MARGIN_BOTTOM})"
    )
    args = parser.parse_args()

    input_files = args.input

    # TỰ ĐỘNG TÌM FILE: Nếu người dùng click đúp (không truyền tham số)
    if not input_files:
        print("Đang quét thư mục hiện tại để tìm các file .pptx...")
        # Lấy tất cả file pptx, loại trừ các file đang mở (bắt đầu bằng ~) và file đã clean
        input_files = [p for p in Path('.').glob('*.pptx') 
                       if not p.stem.endswith('_clean') and not p.name.startswith('~')]

    if not input_files:
        print("❌ Lỗi: Không tìm thấy file .pptx nào cần xử lý trong thư mục hiện tại.", file=sys.stderr)
        input("\nNhấn Enter để thoát...")
        return

    print(f"Đã tìm thấy {len(input_files)} file cần xử lý. Bắt đầu...")

    for input_path in input_files:
        if not input_path.exists():
            print(f"❌ Lỗi: Không tìm thấy file: {input_path}", file=sys.stderr)
            continue

        if input_path.suffix.lower() != ".pptx":
            print(f"⚠️ Bỏ qua file không phải PPTX: {input_path}", file=sys.stderr)
            continue

        # Đặt tên file đầu ra
        if args.output and args.output.name != "clean.pptx" and len(input_files) == 1:
            output_path = args.output
        else:
            output_path = input_path.with_stem(input_path.stem + "_clean")

        try:
            processed, skipped = process_pptx(
                input_path, output_path,
                args.wm_width, args.wm_height,
                args.margin_right, args.margin_bottom,
            )
            status = f"✅ Xong: {input_path.name} -> {output_path.name}"
            status += f" ({processed} ảnh đã xóa watermark"
            if skipped:
                status += f", bỏ qua {skipped} ảnh không phù hợp"
            status += ")"
            print(status)
        except Exception as e:
            print(f"❌ Lỗi khi xử lý file {input_path.name}: {e}", file=sys.stderr)

    print("\n🎉 Hoàn tất toàn bộ!")
    input("Nhấn Enter để thoát...")


if __name__ == "__main__":
    main()
